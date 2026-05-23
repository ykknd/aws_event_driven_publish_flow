from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Inches


def build_pptx_from_manifest(manifest: dict[str, Any], output_path: Path) -> Path:
    output_path.parent.mkdir(parents=True, exist_ok=True)
    presentation = Presentation()

    title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title_slide.shapes.title.text = manifest["job_id"]
    title_slide.placeholders[1].text = manifest["report_type"]

    for artifact in manifest.get("artifacts", []):
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = artifact["label"]
        if artifact["kind"] == "png":
            slide.shapes.add_picture(str(Path(artifact["local_path"])), Inches(1), Inches(1.5), width=Inches(8))
        else:
            textbox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(3))
            textbox.text_frame.text = artifact["s3_key"]

    presentation.save(output_path)
    return output_path

